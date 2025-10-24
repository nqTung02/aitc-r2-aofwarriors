import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import json
import os
import tempfile
import sys

# Redirect print statements to a log file to avoid Unicode errors in the terminal
log_file_path = 'generation.log'
sys.stdout = open(log_file_path, 'w', encoding='utf-8')
sys.stderr = sys.stdout

# --- 0. Cấu hình và Định nghĩa Dữ liệu ---

# Tên file PPTX đầu ra
OUTPUT_PPTX_FILENAME = "Bao_Cao_Tai_Chinh_Doanh_Nghiep.pptx"

# Đường dẫn đến file dữ liệu tài chính
FINANCIAL_DATA_PATH = "data/financial_highlights.json"

# Định nghĩa các slide theo yêu cầu
slide_definitions = [
    {
        "slide_type": "title_slide",
        "title": "Báo Cáo Tài Chính và Hoạt Động Doanh Nghiệp Q2/2024",
        "subtitle": "Phân Tích Hiệu Suất và Chiến Lược Phát Triển",
        "notes": "Trình bày bởi [Tên Người Trình Bày] vào ngày [Ngày]",
        "logo_path": "images/company_logo.png"
    },
    {
        "slide_type": "section_header",
        "title": "I. Tổng Quan Hiệu Suất Tài Chính",
        "subtitle": "Phân tích các chỉ số tài chính cốt lõi"
    },
    {
        "slide_type": "title_and_content",
        "title": "Tóm Tắt Kết Quả Kinh Doanh Q2/2024",
        "content_bullets": [
            "Doanh thu tăng trưởng 15% so với cùng kỳ năm trước, đạt 5.2 triệu USD.",
            "Lợi nhuận ròng đạt 1.8 triệu USD, biên lợi nhuận 34.6%.",
            "Đầu tư vào R&D tăng 20% nhằm thúc đẩy đổi mới sản phẩm."
        ],
        "image_path": "images/q2_summary_icon.png"
    },
    {
        "slide_type": "title_and_chart",
        "title": "Biểu Đồ Doanh Thu Theo Sản Phẩm Q2/2024",
        "chart_definition": {
            "data_source_title": "1H25 Financial Highlights",
            "data_key": "Total operating income",
            "chart_type": "bar",
            "x_axis_keys": ["2Q24", "3Q24", "4Q24", "1Q25", "2Q25"],
            "chart_title": "Tổng Thu Nhập Hoạt Động (Tỷ VND)",
            "x_label": "Quý",
            "y_label": "Thu Nhập (Tỷ VND)",
            "color": "#1f77b4"
        },
        "notes": "Tổng thu nhập hoạt động cho thấy sự biến động qua các quý."
    },
    {
        "slide_type": "title_and_chart",
        "title": "Xu Hướng Lợi Nhuận Trước Thuế",
        "chart_definition": {
            "data_source_title": "1H25 Financial Highlights",
            "data_key": "Profit before tax",
            "chart_type": "line",
            "x_axis_keys": ["2Q24", "3Q24", "4Q24", "1Q25", "2Q25"],
            "chart_title": "Lợi Nhuận Trước Thuế (Tỷ VND)",
            "x_label": "Quý",
            "y_label": "Lợi Nhuận (Tỷ VND)",
            "color": "#d62728"
        },
        "content_bullets": [
            "Lợi nhuận có sự tăng trưởng trở lại trong Q2/2025.",
            "Chi phí hoạt động được kiểm soát hiệu quả."
        ]
    }
]

# Đường dẫn ảnh tĩnh
image_assets = {
    "company_logo": "images/company_logo.png",
    "q2_summary_icon": "images/q2_summary_icon.png",
    "rnd_investment": "images/rnd_investment.jpg",
    "marketing_strategy": "images/marketing_strategy.png",
    "sales_achievement": "images/sales_achievement.png",
    "team_photo": "images/team_photo.jpg"
}

# --- 1. Đọc dữ liệu ---
def load_financial_data(filepath):
    """
    Đọc dữ liệu tài chính từ tệp JSON.
    """
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"Data file not found: {filepath}")
    with open(filepath, 'r', encoding='utf-8') as f:
        data = json.load(f)
    print(f"Successfully loaded data from {filepath}.")
    return data

financial_data = load_financial_data(FINANCIAL_DATA_PATH)

# --- 2. Hàm trợ giúp tạo biểu đồ ---
def create_chart_image(financial_data, chart_definition, output_path):
    """
    Tạo biểu đồ từ dữ liệu tài chính và lưu dưới dạng ảnh.
    Hỗ trợ biểu đồ cột (bar) và đường (line).
    """
    data_source_title = chart_definition["data_source_title"]
    data_key = chart_definition["data_key"]
    chart_type = chart_definition["chart_type"]
    x_axis_keys = chart_definition["x_axis_keys"]
    chart_title = chart_definition["chart_title"]
    x_label = chart_definition.get("x_label", "")
    y_label = chart_definition.get("y_label", "")
    color = chart_definition.get("color", "skyblue")

    # Tìm dữ liệu theo data_source_title và data_key
    # Correcting the logic to handle the list-based JSON structure
    target_data_list = None
    if isinstance(financial_data, dict):
        target_data_list = financial_data.get(data_source_title)
    elif isinstance(financial_data, list) and len(financial_data) > 0:
        # Assuming the first table in the list corresponds to the first chart, etc.
        # This is a fallback if the dictionary structure is not present.
        # A more robust solution might involve inspecting table content.
        if data_source_title == "1H25 Financial Highlights":
             target_data_list = financial_data[0] # First table
        elif len(financial_data) > 1:
             target_data_list = financial_data[1] # Second table as a fallback
    
    if not target_data_list:
        print(f"Warning: Could not find '{data_source_title}' in the financial data.")
        return None

    metric_data = next((item for item in target_data_list if item.get("Balance sheet (VND Bn)") == data_key), None)
    if not metric_data:
        print(f"Warning: Metric '{data_key}' not found in '{data_source_title}'.")
        return None

    values_dict = metric_data
    
    # Chuẩn bị dữ liệu cho biểu đồ
    x_values = x_axis_keys
    y_values = [float(values_dict.get(key, "0").replace(",", "")) for key in x_axis_keys] # Sử dụng 0 nếu không tìm thấy key

    if not any(y_values): # Kiểm tra nếu tất cả y_values đều là 0
        print(f"Warning: No data available to plot the chart for '{data_key}'.")
        return None

    fig, ax = plt.subplots(figsize=(8, 4.5)) # Kích thước hợp lý cho slide

    if chart_type == "bar":
        ax.bar(x_values, y_values, color=color)
    elif chart_type == "line":
        ax.plot(x_values, y_values, marker='o', color=color, linewidth=2)
    else:
        print(f"Warning: Chart type '{chart_type}' is not supported.")
        return None

    ax.set_title(chart_title, fontsize=14)
    ax.set_xlabel(x_label, fontsize=10)
    ax.set_ylabel(y_label, fontsize=10)
    ax.grid(axis='y', linestyle='--', alpha=0.7)
    plt.xticks(rotation=45, ha='right')
    plt.tight_layout()

    try:
        plt.savefig(output_path, dpi=300)
        print(f"Created chart '{chart_title}' and saved to {output_path}")
        return output_path
    except Exception as e:
        print(f"Error saving chart to {output_path}: {e}")
        return None
    finally:
        plt.close(fig) # Đóng figure để giải phóng bộ nhớ

# --- 3. Hàm trợ giúp tạo slide ---

def add_title_slide(prs, slide_def):
    """Thêm slide tiêu đề."""
    slide_layout = prs.slide_layouts[0] # Bố cục slide tiêu đề
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1] # Placeholder subtitle thường là index 1

    title.text = slide_def["title"]
    subtitle.text = slide_def["subtitle"]

    # Thêm ghi chú dưới dạng hộp văn bản mới nếu không có placeholder phù hợp
    left = Inches(1)
    top = Inches(6)
    width = Inches(8)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = slide_def["notes"]
    p.font.size = Pt(12)
    p.font.italic = True

    # Chèn logo công ty
    if slide_def.get("logo_path") and os.path.exists(slide_def["logo_path"]):
        left = Inches(8)
        top = Inches(0.5)
        height = Inches(1.0) # Điều chỉnh chiều cao logo
        slide.shapes.add_picture(slide_def["logo_path"], left, top, height=height)
    print(f"Created title slide: '{slide_def['title']}'")

def add_section_header_slide(prs, slide_def):
    """Thêm slide tiêu đề phần."""
    slide_layout = prs.slide_layouts[5] # Bố cục chỉ có tiêu đề
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = slide_def["title"]
    
    # Thêm subtitle dưới dạng hộp văn bản mới
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = slide_def["subtitle"]
    p.font.size = Pt(24)
    p.font.bold = True
    print(f"Created section header slide: '{slide_def['title']}'")

def add_title_and_content_slide(prs, slide_def):
    """Thêm slide tiêu đề và nội dung (có thể kèm ảnh)."""
    slide_layout = prs.slide_layouts[1] # Bố cục tiêu đề và nội dung
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = slide_def["title"]

    # Placeholder cho nội dung chính
    body_shape = slide.placeholders[1] 
    tf = body_shape.text_frame
    tf.clear() # Xóa nội dung mặc định

    # Thêm các bullet points
    if slide_def.get("content_bullets"):
        for bullet_text in slide_def["content_bullets"]:
            p = tf.add_paragraph()
            p.text = bullet_text
            p.level = 1 # Cấp độ bullet point
            p.font.size = Pt(18)
    
    # Chèn ảnh nếu có
    if slide_def.get("image_path") and os.path.exists(slide_def["image_path"]):
        # Vị trí ảnh: bên phải của nội dung
        img_left = Inches(6.5)
        img_top = Inches(2.5)
        img_height = Inches(3.5)
        slide.shapes.add_picture(slide_def["image_path"], img_left, img_top, height=img_height)
        # Điều chỉnh kích thước và vị trí của placeholder nội dung để không chồng chéo ảnh
        body_shape.left = Inches(0.5)
        body_shape.top = Inches(1.8)
        body_shape.width = Inches(5.5)
        body_shape.height = Inches(4.5)
        
    print(f"Created title and content slide: '{slide_def['title']}'")

def add_title_and_chart_slide(prs, slide_def, financial_data, temp_chart_dir):
    """Thêm slide tiêu đề và biểu đồ (có thể kèm bullet points)."""
    slide_layout = prs.slide_layouts[1] # Bố cục tiêu đề và nội dung
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = slide_def["title"]

    chart_def = slide_def["chart_definition"]
    chart_filename = f"chart_{chart_def['data_key'].replace(' ', '_')}_{chart_def['chart_type']}.png"
    chart_path = os.path.join(temp_chart_dir, chart_filename)

    # Tạo biểu đồ và nhúng vào slide
    created_chart_path = create_chart_image(financial_data, chart_def, chart_path)
    if created_chart_path:
        # Vị trí cho biểu đồ (nửa bên trái)
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(6)
        slide.shapes.add_picture(created_chart_path, left, top, width=width)
        
        # Thêm bullet points nếu có (nửa bên phải)
        if slide_def.get("content_bullets"):
            left = Inches(6.8)
            top = Inches(2.0)
            width = Inches(3.5)
            height = Inches(4.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.clear()
            tf.word_wrap = True # Đảm bảo văn bản xuống dòng tự động
            
            for bullet_text in slide_def["content_bullets"]:
                p = tf.add_paragraph()
                p.text = bullet_text
                p.level = 1
                p.font.size = Pt(16)
                p.space_after = Pt(10) # Khoảng cách giữa các bullet
                
            # Đặt auto_size để hộp văn bản tự điều chỉnh theo nội dung
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
    # Thêm ghi chú nếu có
    if slide_def.get("notes"):
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(9)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = slide_def["notes"]
        p.font.size = Pt(12)
        p.font.italic = True

    print(f"Created chart slide: '{slide_def['title']}'")

# --- 4. Logic chính để tạo Presentation ---

def create_presentation(slide_defs, financial_data, output_filename):
    """
    Tạo một presentation PowerPoint hoàn chỉnh dựa trên các định nghĩa slide.
    """
    prs = Presentation()

    # Tạo thư mục tạm thời để lưu các biểu đồ
    with tempfile.TemporaryDirectory() as temp_chart_dir:
        print(f"Temporary directory for charts: {temp_chart_dir}")

        for slide_def in slide_defs:
            slide_type = slide_def["slide_type"]
            if slide_type == "title_slide":
                add_title_slide(prs, slide_def)
            elif slide_type == "section_header":
                add_section_header_slide(prs, slide_def)
            elif slide_type == "title_and_content":
                add_title_and_content_slide(prs, slide_def)
            elif slide_type == "title_and_chart":
                add_title_and_chart_slide(prs, slide_def, financial_data, temp_chart_dir)
            else:
                print(f"Undefined slide type: {slide_type}. Skipping this slide.")
        
        # Lưu presentation
        prs.save(output_filename)
        print(f"\nPresentation created successfully: {output_filename}")

# --- Chạy chương trình ---
if __name__ == "__main__":
    create_presentation(slide_definitions, financial_data, OUTPUT_PPTX_FILENAME)
